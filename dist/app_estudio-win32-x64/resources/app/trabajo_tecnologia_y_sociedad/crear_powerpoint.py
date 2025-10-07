"""
Script para generar PowerPoint usando python-pptx
Ejecutar: pip install python-pptx
Luego ejecutar este script para generar el archivo .pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

def create_presentation():
    # Crear presentación
    prs = Presentation()
    
    # Configurar tamaño de slide (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Portada
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Título principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "IA, Prevención de Fraude y Trabajo en Fintech"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Una perspectiva socio-ética"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Autores
    authors_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    authors_frame = authors_box.text_frame
    authors_frame.text = "Marto, Rocío y Naty"
    authors_para = authors_frame.paragraphs[0]
    authors_para.font.size = Pt(20)
    authors_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Introducción - Marto
    slide = prs.slides.add_slide(slide_layout)
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Introducción"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Presentador
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "📢 Marto"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """Presentamos nuestro trabajo:

"IA, prevención de fraude y trabajo en Fintech: 
una perspectiva socio-ética"

Equipo: Marto, Rocío y Naty"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(20)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 3: Objetivo - Marto
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Objetivo Principal"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "🎯 Marto"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """Analizar cómo se aplica la inteligencia artificial 
en la detección de fraudes dentro del ecosistema Fintech

Enfocándonos en:
• Aspectos técnicos
• Aspectos éticos  
• Aspectos laborales
• Aspectos sociales"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 4: Hipótesis - Marto
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Hipótesis Central"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "💡 Marto"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """"La inteligencia artificial no es neutral"

Su diseño, entrenamiento y uso tienen impactos directos 
sobre las personas y la sociedad

Por eso, creemos que su implementación debe ser 
crítica, responsable y ética"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 5: Deepfakes - Rocío
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Deepfakes: La Nueva Amenaza"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "🧠 Rocío"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """Tecnología basada en IA generativa que permite imitar:

🎭 Rostros de forma hiperrealista
🎤 Voces sintéticas  
📄 Documentos falsificados

Representa una amenaza directa al ecosistema Fintech"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 6: Impacto en Fintech - Rocío
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Impacto en el Entorno Fintech"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "⚠️ Rocío"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """Los deepfakes pueden:

• Burlar sistemas de verificación biométrica
• Comprometer el reconocimiento facial
• Vulnerar la verificación de voz
• Generar fraudes sofisticados

Esto representa una amenaza directa para la seguridad financiera"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 7: Casos Reales - Rocío
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Casos de Estudio"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "📊 Rocío"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """Casos reales analizados:

• Personal Pay: Intentos de fraude con identidades sintéticas

• PAYGILANT: Implementación de herramientas de detección

• Klarna: Reducción del 40% de personal tras implementar IA"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 8: Entrevista - Rocío
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Perspectiva desde la Industria"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "🎤 Rocío"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """Entrevista con Vanesa Fabres
Analista Senior de Fraudes en Moni Online

Temas abordados:
• Uso actual de herramientas de IA
• Desafíos: falsos positivos
• Regulaciones vigentes
• Límites de la automatización"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 9: Marco Ético - Naty
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Perspectiva Ética"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "📚 Naty"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """Marco teórico:

Gabriela Ramos (UNESCO)
"La IA mal gobernada puede profundizar desigualdades"

Enrique Marí  
"La tecnología no es neutral y siempre conlleva una carga aplicativa"

Ambos autores coinciden en la importancia del diseño ético"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 10: Diseño Ético - Naty
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Diseño Ético"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "⚖️ Naty"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """Principio fundamental:

"El diseño ético debe estar desde el inicio 
y no ser algo agregado después"

La ética debe ser inherente al proceso de desarrollo,
no un complemento posterior"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 11: Conclusión - Naty
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusión"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "🎯 Naty"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = """El futuro de la prevención de fraude en Fintech:

ALIANZA ESTRATÉGICA
entre inteligencia artificial avanzada 
y criterio humano experto

Con una gobernanza que priorice:
• La justicia
• La transparencia  
• La protección de derechos"""
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 12: Gracias
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Muchas Gracias"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "¿Preguntas?"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    authors_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    authors_frame = authors_box.text_frame
    authors_frame.text = "Marto, Rocío y Naty - Trabajo de Tecnología y Sociedad"
    authors_para = authors_frame.paragraphs[0]
    authors_para.font.size = Pt(16)
    authors_para.alignment = PP_ALIGN.CENTER
    
    # Guardar presentación
    prs.save('presentacion_IA_fraude_fintech.pptx')
    print("Presentación creada exitosamente: presentacion_IA_fraude_fintech.pptx")

if __name__ == "__main__":
    create_presentation()
