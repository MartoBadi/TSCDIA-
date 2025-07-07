"""
Script para generar PowerPoint con tema oscuro usando python-pptx
Ejecutar: pip install python-pptx
Luego ejecutar este script para generar el archivo .pptx con colores invertidos
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR

def create_dark_presentation():
    # Crear presentación
    prs = Presentation()
    
    # Configurar tamaño de slide (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Colores del tema oscuro
    DARK_BG = RGBColor(26, 26, 46)  # #1a1a2e
    WHITE_TEXT = RGBColor(224, 224, 224)  # #e0e0e0
    GOLD_ACCENT = RGBColor(255, 215, 0)  # #ffd700
    LIGHT_GRAY = RGBColor(240, 240, 240)  # #f0f0f0
    
    # Slide 1: Portada
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo oscuro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Título principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "IA, Prevención de Fraude y Trabajo en Fintech"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Una perspectiva socio-ética - Tema Oscuro"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.color.rgb = GOLD_ACCENT
    
    # Autores
    authors_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    authors_frame = authors_box.text_frame
    authors_frame.text = "Marto, Rocío y Naty"
    authors_para = authors_frame.paragraphs[0]
    authors_para.font.size = Pt(20)
    authors_para.alignment = PP_ALIGN.CENTER
    authors_para.font.color.rgb = WHITE_TEXT
    
    # Función para crear slide con fondo oscuro
    def create_dark_slide(title_text, presenter_text, content_text):
        slide = prs.slides.add_slide(slide_layout)
        
        # Fondo oscuro
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        # Título
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Presentador
        presenter_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.5), Inches(3), Inches(0.5))
        presenter_frame = presenter_box.text_frame
        presenter_frame.text = presenter_text
        presenter_para = presenter_frame.paragraphs[0]
        presenter_para.font.size = Pt(16)
        presenter_para.alignment = PP_ALIGN.RIGHT
        presenter_para.font.color.rgb = GOLD_ACCENT
        
        # Contenido
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
        content_frame = content_box.text_frame
        content_frame.text = content_text
        
        for para in content_frame.paragraphs:
            para.font.size = Pt(18)
            para.alignment = PP_ALIGN.CENTER
            para.font.color.rgb = WHITE_TEXT
            
        return slide
    
    # Slide 2: Introducción - Marto
    create_dark_slide(
        "Introducción",
        "📢 Marto",
        """Presentamos nuestro trabajo:

"IA, prevención de fraude y trabajo en Fintech: 
una perspectiva socio-ética"

Equipo: Marto, Rocío y Naty"""
    )
    
    # Slide 3: Objetivo - Marto
    create_dark_slide(
        "Objetivo Principal",
        "🎯 Marto",
        """Analizar cómo se aplica la inteligencia artificial 
en la detección de fraudes dentro del ecosistema Fintech

Enfocándonos en:
• Aspectos técnicos
• Aspectos éticos  
• Aspectos laborales
• Aspectos sociales"""
    )
    
    # Slide 4: Hipótesis - Marto
    create_dark_slide(
        "Hipótesis Central",
        "💡 Marto",
        """"La inteligencia artificial no es neutral"

Su diseño, entrenamiento y uso tienen impactos directos 
sobre las personas y la sociedad

Por eso, creemos que su implementación debe ser 
crítica, responsable y ética"""
    )
    
    # Slide 5: Deepfakes - Rocío
    create_dark_slide(
        "Deepfakes: La Nueva Amenaza",
        "🧠 Rocío",
        """Tecnología basada en IA generativa que permite imitar:

🎭 Rostros de forma hiperrealista
🎤 Voces sintéticas  
📄 Documentos falsificados

Representa una amenaza directa al ecosistema Fintech"""
    )
    
    # Slide 6: Impacto en Fintech - Rocío
    create_dark_slide(
        "Impacto en el Entorno Fintech",
        "⚠️ Rocío",
        """Los deepfakes pueden:

• Burlar sistemas de verificación biométrica
• Comprometer el reconocimiento facial
• Vulnerar la verificación de voz
• Generar fraudes sofisticados

Esto representa una amenaza directa para la seguridad financiera"""
    )
    
    # Slide 7: Casos Reales - Rocío
    create_dark_slide(
        "Casos de Estudio",
        "📊 Rocío",
        """Casos reales analizados:

• Personal Pay: Intentos de fraude con identidades sintéticas

• PAYGILANT: Implementación de herramientas de detección

• Klarna: Reducción del 40% de personal tras implementar IA"""
    )
    
    # Slide 8: Entrevista - Rocío
    create_dark_slide(
        "Perspectiva desde la Industria",
        "🎤 Rocío",
        """Entrevista con Vanesa Fabres
Analista Senior de Fraudes en Moni Online

Temas abordados:
• Uso actual de herramientas de IA
• Desafíos: falsos positivos
• Regulaciones vigentes
• Límites de la automatización"""
    )
    
    # Slide 9: Marco Ético - Naty
    create_dark_slide(
        "Perspectiva Ética",
        "📚 Naty",
        """Marco teórico:

Gabriela Ramos (UNESCO)
"La IA mal gobernada puede profundizar desigualdades"

Enrique Marí  
"La tecnología no es neutral y siempre conlleva una carga aplicativa"

Ambos autores coinciden en la importancia del diseño ético"""
    )
    
    # Slide 10: Diseño Ético - Naty
    create_dark_slide(
        "Diseño Ético",
        "⚖️ Naty",
        """Principio fundamental:

"El diseño ético debe estar desde el inicio 
y no ser algo agregado después"

La ética debe ser inherente al proceso de desarrollo,
no un complemento posterior"""
    )
    
    # Slide 11: Conclusión - Naty
    create_dark_slide(
        "Conclusión",
        "🎯 Naty",
        """El futuro de la prevención de fraude en Fintech:

ALIANZA ESTRATÉGICA
entre inteligencia artificial avanzada 
y criterio humano experto

Con una gobernanza que priorice:
• La justicia
• La transparencia  
• La protección de derechos"""
    )
    
    # Slide 12: Gracias
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo oscuro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Muchas Gracias"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "¿Preguntas?"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.color.rgb = GOLD_ACCENT
    
    authors_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    authors_frame = authors_box.text_frame
    authors_frame.text = "Marto, Rocío y Naty - Trabajo de Tecnología y Sociedad"
    authors_para = authors_frame.paragraphs[0]
    authors_para.font.size = Pt(16)
    authors_para.alignment = PP_ALIGN.CENTER
    authors_para.font.color.rgb = WHITE_TEXT
    
    # Guardar presentación
    prs.save('presentacion_IA_fraude_fintech_dark.pptx')
    print("Presentación con tema oscuro creada exitosamente: presentacion_IA_fraude_fintech_dark.pptx")

if __name__ == "__main__":
    create_dark_presentation()
